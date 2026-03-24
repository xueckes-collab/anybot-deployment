"""
文档提取服务：从各种格式（PDF、DOCX、XLSX、PPTX 等）提取文本
并将其添加到知识库
"""

import sys
import json
from pathlib import Path
from datetime import datetime

# 文档处理库
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from config import get_settings

# 后台管理的文档目录
ADMIN_DOCS_DIR = Path(__file__).resolve().parent.parent.parent / "admin" / "admin-data" / "documents"
# 提取后的文档目录
EXTRACTED_DIR = Path(__file__).resolve().parent.parent / "data" / "extracted"
EXTRACTED_DIR.mkdir(parents=True, exist_ok=True)

# 目录文件
CATALOG_FILE = EXTRACTED_DIR / "catalog.json"


def extract_pdf(filepath: str) -> str:
    """从 PDF 提取文本"""
    try:
        text = []
        with open(filepath, 'rb') as f:
            reader = PdfReader(f)
            for page in reader.pages:
                text.append(page.extract_text())
        return "\n".join(text)
    except Exception as e:
        print(f"  ❌ PDF 提取失败: {e}")
        return ""


def extract_docx(filepath: str) -> str:
    """从 DOCX 提取文本"""
    try:
        doc = DocxDocument(filepath)
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text.append(cell.text)
        return "\n".join(text)
    except Exception as e:
        print(f"  ❌ DOCX 提取失败: {e}")
        return ""


def extract_doc(filepath: str) -> str:
    """从 DOC 提取文本（使用 python-docx，仅支持 .docx）"""
    # 旧版 .doc 格式需要特殊处理，这里尝试用 docx 库
    try:
        doc = DocxDocument(filepath)
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        return "\n".join(text)
    except Exception as e:
        print(f"  ⚠️  DOC 提取失败（仅支持 .docx）: {e}")
        return ""


def extract_xlsx(filepath: str) -> str:
    """从 XLSX 提取文本"""
    try:
        wb = load_workbook(filepath)
        text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text.append(f"【Sheet: {sheet}】")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    text.append(row_text)
        return "\n".join(text)
    except Exception as e:
        print(f"  ❌ XLSX 提取失败: {e}")
        return ""


def extract_xls(filepath: str) -> str:
    """从 XLS 提取文本（使用 openpyxl，仅支持 .xlsx）"""
    try:
        wb = load_workbook(filepath)
        text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text.append(f"【Sheet: {sheet}】")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    text.append(row_text)
        return "\n".join(text)
    except Exception as e:
        print(f"  ⚠️  XLS 提取失败（仅支持 .xlsx）: {e}")
        return ""


def extract_pptx(filepath: str) -> str:
    """从 PPTX 提取文本"""
    try:
        prs = Presentation(filepath)
        text = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            text.append(f"【Slide {slide_idx}】")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"  ❌ PPTX 提取失败: {e}")
        return ""


def extract_ppt(filepath: str) -> str:
    """从 PPT 提取文本（使用 python-pptx，仅支持 .pptx）"""
    try:
        prs = Presentation(filepath)
        text = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            text.append(f"【Slide {slide_idx}】")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"  ⚠️  PPT 提取失败（仅支持 .pptx）: {e}")
        return ""


def extract_txt(filepath: str) -> str:
    """从 TXT 读取文本"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"  ❌ TXT 读取失败: {e}")
        return ""


def extract_md(filepath: str) -> str:
    """从 Markdown 读取文本"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"  ❌ Markdown 读取失败: {e}")
        return ""


def extract_csv(filepath: str) -> str:
    """从 CSV 读取文本"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"  ❌ CSV 读取失败: {e}")
        return ""


def extract_json(filepath: str) -> str:
    """从 JSON 读取文本"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            data = json.load(f)
            return json.dumps(data, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"  ❌ JSON 读取失败: {e}")
        return ""


def extract_xml(filepath: str) -> str:
    """从 XML 读取文本"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"  ❌ XML 读取失败: {e}")
        return ""


def extract_html(filepath: str) -> str:
    """从 HTML 读取文本"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"  ❌ HTML 读取失败: {e}")
        return ""


def extract_rtf(filepath: str) -> str:
    """从 RTF 读取文本（简单提取）"""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
            # 简单的 RTF 文本提取
            import re
            text = re.sub(r'\\[a-z]+\d*\s?', '', content)
            text = re.sub(r'[{}]', '', text)
            return text
    except Exception as e:
        print(f"  ❌ RTF 读取失败: {e}")
        return ""


EXTRACTORS = {
    '.pdf': extract_pdf,
    '.txt': extract_txt,
    '.docx': extract_docx,
    '.doc': extract_doc,
    '.xlsx': extract_xlsx,
    '.xls': extract_xls,
    '.pptx': extract_pptx,
    '.ppt': extract_ppt,
    '.md': extract_md,
    '.csv': extract_csv,
    '.json': extract_json,
    '.xml': extract_xml,
    '.html': extract_html,
    '.htm': extract_html,
    '.rtf': extract_rtf,
}


def load_catalog() -> dict:
    """加载目录"""
    if CATALOG_FILE.exists():
        return json.loads(CATALOG_FILE.read_text(encoding='utf-8'))
    return {}


def save_catalog(catalog: dict):
    """保存目录"""
    CATALOG_FILE.write_text(json.dumps(catalog, ensure_ascii=False, indent=2), encoding='utf-8')


def process_documents():
    """处理后台上传的所有文档"""
    if not ADMIN_DOCS_DIR.exists():
        print(f"❌ 文档目录不存在: {ADMIN_DOCS_DIR}")
        return

    catalog = load_catalog()
    processed = 0

    print(f"\n📂 扫描文档目录: {ADMIN_DOCS_DIR}")
    print(f"📚 提取目录: {EXTRACTED_DIR}\n")

    for doc_file in ADMIN_DOCS_DIR.iterdir():
        if not doc_file.is_file():
            continue

        ext = doc_file.suffix.lower()
        if ext not in EXTRACTORS:
            print(f"⏭️  跳过不支持的格式: {doc_file.name}")
            continue

        # 检查是否已处理
        doc_key = doc_file.name
        if doc_key in catalog:
            print(f"✅ 已处理: {doc_file.name}")
            continue

        print(f"🔄 处理: {doc_file.name}")

        # 提取文本
        extractor = EXTRACTORS[ext]
        content = extractor(str(doc_file))

        if not content or not content.strip():
            print(f"  ⚠️  无法提取内容")
            continue

        # 保存提取的文本
        output_file = EXTRACTED_DIR / f"{doc_file.stem}.md"
        output_file.write_text(content, encoding='utf-8')
        print(f"  ✅ 已保存: {output_file.name}")

        # 更新目录
        catalog[doc_key] = {
            "file": f"{doc_file.stem}.md",
            "name": doc_file.stem,
            "type": ext[1:].upper(),  # 去掉点
            "source": doc_file.name,
            "extracted_at": datetime.now().isoformat(),
        }
        processed += 1

    save_catalog(catalog)
    print(f"\n✅ 处理完成！共处理 {processed} 个文档")
    print(f"📊 目录中共有 {len(catalog)} 个文档")


if __name__ == "__main__":
    process_documents()
